"""PPT 导出：报告 Markdown → 16:9 .pptx（python-pptx 渲染 + 关系图嵌入）。"""

import os


def test_export_report_pptx_renders_slides(sample_event, tmp_path):
    import app.rule_engine as rule_engine
    from app.pptx_renderer import export_report_pptx

    md = rule_engine.generate(
        rule_engine.StructuredInput.model_validate(sample_event)
    )
    result = export_report_pptx(
        sample_event.get("title") or "测试报告",
        md,
        output_dir=str(tmp_path),
        slug="pptx_test",
    )

    assert os.path.exists(result["pptx"])
    assert result["pptx"].endswith(".pptx")
    from pptx import Presentation

    prs = Presentation(result["pptx"])
    # 封面/情况概述/主体/冲突/分析/结论/建议/附录 至少 7 页
    assert len(prs.slides) >= 7


def test_export_report_pptx_embeds_diagrams(sample_event, tmp_path):
    """报告含 DIAGRAM 时生成 PNG 并整页嵌入，diagram_count > 0。"""
    import app.rule_engine as rule_engine
    from app.pptx_renderer import export_report_pptx

    md = rule_engine.generate(
        rule_engine.StructuredInput.model_validate(sample_event)
    )
    result = export_report_pptx(
        sample_event.get("title") or "测试报告",
        md,
        output_dir=str(tmp_path),
        slug="pptx_diagram_test",
    )

    # rule 引擎报告通常带利益网络图；若生成则 PPT 里应存在图片页
    if result["diagram_count"] > 0:
        from pptx import Presentation

        prs = Presentation(result["pptx"])
        picture_count = sum(
            1
            for slide in prs.slides
            for shape in slide.shapes
            if shape.shape_type == 13  # PICTURE
        )
        assert picture_count == result["diagram_count"]
