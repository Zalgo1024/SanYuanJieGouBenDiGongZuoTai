"""PPT 渲染器 — 从报告 Markdown 生成 16:9 演示稿（.pptx）。

v2 版式引擎（方案 A）：
- 建立设计系统：12 列网格、字号阶梯、色彩角色、统一间距；
- 用「页面原型」替代过去的「标题 + 一段文字」单一骨架：封面 / 章节页 /
  叙述页 / 冲突卡 / 主体表 / 六类利益矩阵 / 关系网络 / 全幅关系图 /
  结论引言 / 编号建议 / 来源双栏；
- 补齐 v1 的重大缺陷：v1 只解析 subheading/paragraph/list，
  **表格（table）被整块丢弃**，导致主体清单、六类利益分布、关系网络
  这些结构化程度最高的内容在 PPT 里是空的。v2 全部渲染为真正的表格/矩阵。

入口：export_report_pptx(title, markdown, output_dir, slug)
      -> {"pptx", "folder", "diagram_count", "diagrams"}
"""
import os
import re
import sys
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from app.settings import settings

# ── 三元结构六类利益配色（与 theory_config.json 一致）──
C = {
    "gold": RGBColor(0xC9, 0xA2, 0x27),     # 物金
    "green": RGBColor(0x3E, 0x8E, 0x5A),    # 安绿
    "purple": RGBColor(0x5B, 0x4B, 0xA8),   # 政紫
    "orange": RGBColor(0xD9, 0x7A, 0x2B),   # 身份橙
    "blue": RGBColor(0x2F, 0x6F, 0xC0),     # 制度蓝
    "cyan": RGBColor(0x17, 0x9B, 0x9B),     # 公共青
}
INK = RGBColor(0x1A, 0x22, 0x33)
MUTED = RGBColor(0x62, 0x6C, 0x7E)
FAINT = RGBColor(0x9A, 0xA3, 0xB2)
LINE = RGBColor(0xE2, 0xE7, 0xEF)
SURFACE = RGBColor(0xF6, 0xF8, 0xFC)
CANVAS = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xB5, 0x42, 0x38)
ACCENT = C["purple"]

# 六类利益表头 → 配色（按关键词匹配）
INTEREST_HEADER_COLORS = [
    ("物质", C["gold"]),
    ("安全", C["green"]),
    ("政治", C["purple"]),
    ("身份", C["orange"]),
    ("制度", C["blue"]),
    ("公共", C["cyan"]),
]

FONT = "微软雅黑"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── 网格系统：12 列，左右边距 0.7" ──
MARGIN = Inches(0.7)
CONTENT_W = SLIDE_W - 2 * MARGIN
GRID = 12
GUTTER = Inches(0.22)
COL_W = Emu(int((CONTENT_W - (GRID - 1) * GUTTER) / GRID))

# 版心纵向节奏
TITLE_TOP = Inches(0.46)
TITLE_H = Inches(1.02)
BODY_TOP = Inches(1.78)
BODY_H = Inches(5.05)
FOOTER_TOP = Inches(6.92)

LINK_RE = re.compile(r"\[([^\]]+)\]\((https?://[^)\s]+)\)")
CONFLICT_RE = re.compile(
    r"[【\[]([^】\]]+)[】\]]\s*对\s*[【\[]([^】\]]+)[】\]]\s*在\s*[「\"']?([^」\"']*)[」\"']?\s*上的\s*(\S+?)\s*[:：]"
)
SOURCE_RE = re.compile(r"[（(]\s*来源[:：]\s*([^）)]+)[）)]")
BRACKET_RE = re.compile(r"[【\[]([^】\]]+)[】\]]")


# ══════════════════════════════════════════════════════════════
# 基础设施
# ══════════════════════════════════════════════════════════════
def _ensure_engine_path() -> None:
    d = settings.engine_dir
    if d and os.path.isdir(d) and d not in sys.path:
        sys.path.insert(0, d)


def col_left(i: int) -> int:
    """第 i 列（0 起）的左边缘。"""
    return Emu(int(MARGIN) + i * int(COL_W + GUTTER))


def col_span(n: int) -> int:
    """跨 n 列的总宽度。"""
    return Emu(n * int(COL_W) + (n - 1) * int(GUTTER))


def _clean(text: str) -> str:
    """去掉 Markdown 链接语法与多余空白。"""
    text = LINK_RE.sub(r"\1", text or "")
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    return text.strip()


def _apply_font(run, size, color=INK, bold=False) -> None:
    """设置 run 字体（含中文字体 east asian）。"""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def _para(tf, text, size, color=INK, bold=False, first=False,
          space_before=0, space_after=0, align=PP_ALIGN.LEFT, line_spacing=1.35):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if text:
        run = p.add_run()
        run.text = text
        _apply_font(run, size, color, bold)
    return p


def _rect(slide, left, top, width, height, fill=None, line=None,
          radius=None, line_width=0.75):
    """矩形 / 圆角矩形。radius 为圆角比例（0.0-0.5）。"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:  # noqa: BLE001
            pass
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.word_wrap = True
    return shape


def _visual_len(text: str) -> float:
    """估算文本视觉宽度（单位 em）：中日韩全角记 1.0，西文半角记 0.56。"""
    total = 0.0
    for ch in text or "":
        total += 1.0 if ord(ch) > 0x2E80 else 0.56
    return total


def _pill(slide, left, top, height, text, fill, text_color=CANVAS, size=9.5, bold=True,
          max_width=None):
    """胶囊徽章：按真实字宽自适应；超出 max_width 时自动降字号，保证单行不换行。"""
    pad = Inches(0.15)
    s = float(size)
    vlen = _visual_len(text)
    # 1.08 缓冲：抵消字体度量差异，避免文字刚好卡满而被迫换行
    text_w = int(Inches(vlen * s * 1.08 / 72.0))
    if max_width is not None:
        inner_max = int(max_width) - 2 * int(pad)
        if text_w > inner_max and vlen > 0:
            s = max(7.5, s * (inner_max / float(text_w)))
            text_w = int(Inches(vlen * s * 1.08 / 72.0))
    w = Emu(text_w + 2 * int(pad))
    if max_width is not None:
        w = Emu(min(int(w), int(max_width)))
    shape = _rect(slide, left, top, w, height, fill=fill, radius=0.42)
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, text, s, text_color, bold, first=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
    return shape


def _num_circle(slide, left, top, size, number, fill, text_color=CANVAS):
    """编号圆。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, str(number), 13, text_color, True, first=True,
          align=PP_ALIGN.CENTER, line_spacing=1.0)
    return shape


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title_block(slide, text, eyebrow=None, accent=ACCENT):
    """统一标题区： eyebrow（小标签）+ 大标题 + 强调短横。"""
    top = TITLE_TOP
    if eyebrow:
        tf = _textbox(slide, MARGIN, top, CONTENT_W, Inches(0.26))
        _para(tf, eyebrow.upper(), 10.5, accent, True, first=True)
        top = top + Inches(0.34)
    tf = _textbox(slide, MARGIN, top, CONTENT_W, Inches(0.62))
    _para(tf, text, 27, INK, True, first=True, line_spacing=1.1)
    bar = _rect(slide, MARGIN, top + Inches(0.66), Inches(0.9), Pt(4.5), fill=accent)
    return top + Inches(0.66)


def _footer(slide, text, number=None):
    """页脚：细分隔线 + 左说明 + 右页码。"""
    _rect(slide, MARGIN, FOOTER_TOP, CONTENT_W, Pt(0.75), fill=LINE)
    tf = _textbox(slide, MARGIN, FOOTER_TOP + Inches(0.1), CONTENT_W, Inches(0.28))
    _para(tf, text, 9, FAINT, first=True)
    if number is not None:
        tf2 = _textbox(slide, MARGIN, FOOTER_TOP + Inches(0.1), CONTENT_W, Inches(0.28))
        _para(tf2, str(number).zfill(2), 9.5, FAINT, True, first=True, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# 内容提取
# ══════════════════════════════════════════════════════════════
def _section(report, sid):
    return report.sections.get(sid)


def _items(report, sid, li_limit=10, kinds=("sub", "p", "li")) -> list:
    """提取章节内容为 [(kind, text)]。"""
    sec = _section(report, sid)
    out: list = []
    if not sec:
        return out
    for b in sec.blocks:
        if b.type == "subheading" and b.text.strip() and "sub" in kinds:
            out.append(("sub", b.text.strip()))
        elif b.type in ("paragraph", "quote") and b.text.strip() and "p" in kinds:
            out.append(("p", b.text.strip()))
        elif b.type == "list" and b.items and "li" in kinds:
            for it in b.items[:li_limit]:
                clean = _clean(it)
                if clean:
                    out.append(("li", clean))
    return out


def _tables(report, sid) -> list:
    """提取章节中的表格块：[{caption, rows}]。v1 完全丢弃了表格，v2 补上。"""
    sec = _section(report, sid)
    out: list = []
    if not sec:
        return out
    caption = ""
    for b in sec.blocks:
        if b.type == "subheading":
            caption = b.text.strip()
        elif b.type == "table" and b.rows:
            rows = [[_clean(c) for c in row] for row in b.rows if any(c.strip() for c in row)]
            if rows:
                out.append({"caption": caption, "rows": rows})
            caption = ""
    return out


def _find_table(tables, keywords) -> dict | None:
    """按表头关键词找表格。"""
    for t in tables:
        header = " ".join(t["rows"][0]) if t["rows"] else ""
        if any(k in header for k in keywords):
            return t
    return None


def _parse_conflict(text: str) -> dict:
    """把冲突条目解析为结构化字段。

    输入形如：1. **【A】对【B】在"X"上的张力**：说明...（来源：xxx）
    """
    raw = _clean(text)
    raw = re.sub(r"^\d+[\.、]\s*", "", raw)
    m = CONFLICT_RE.search(raw)
    if not m:
        return {"a": "", "b": "", "interest": "", "tension": "张力", "detail": raw, "source": ""}
    a, b, interest, tension = m.group(1), m.group(2), m.group(3), m.group(4)
    detail = raw[m.end():].strip()
    src = SOURCE_RE.search(detail)
    source = src.group(1).strip() if src else ""
    if src:
        detail = (detail[:src.start()] + detail[src.end():]).strip()
    detail = re.sub(r"^[:：]\s*", "", detail)
    return {
        "a": a.strip(), "b": b.strip(), "interest": interest.strip(),
        "tension": tension.strip(), "detail": detail, "source": source,
    }


def _interest_color(header: str) -> RGBColor:
    for key, color in INTEREST_HEADER_COLORS:
        if key in header:
            return color
    return ACCENT


def _level_color(text: str) -> RGBColor:
    """按强度等级给单元格文字上色。"""
    t = text or ""
    if "受损" in t:
        return DANGER
    if "极高" in t:
        return C["orange"]
    if "极低" in t:
        return FAINT
    if "高" in t:
        return C["purple"]
    if "低" in t:
        return FAINT
    if "中" in t:
        return MUTED
    return MUTED


def _sources(report) -> list:
    """附录来源：[(名称, url)]。"""
    sec = _section(report, "appendix")
    srcs: list = []
    if not sec:
        return srcs
    for b in sec.blocks:
        if b.type == "list" and b.items:
            for it in b.items:
                m = LINK_RE.search(it)
                if m:
                    srcs.append((m.group(1).strip(), m.group(2).strip()))
                else:
                    clean = it.strip()
                    if clean:
                        srcs.append((clean, ""))
    return srcs[:12]


def _diagram_pngs(report, folder: str) -> list:
    """把报告里的 DIAGRAM 块渲染为 PNG（viz_network）。"""
    _ensure_engine_path()
    from viz_network import generate_diagram

    out: list = []
    idx = 0
    for sec in report.section_seq:
        for b in sec.blocks:
            if b.type == "diagram" and b.diagram_data:
                idx += 1
                png = os.path.join(folder, f"diagram_{idx}.png")
                try:
                    generate_diagram(b.diagram_data, png)
                    if os.path.exists(png):
                        out.append({
                            "seq": idx, "png": png,
                            "title": b.diagram_data.get("title") or f"图 {idx}",
                        })
                except Exception as e:  # noqa: BLE001
                    print(f"[pptx] 图 {idx} 生成失败：{e}")
    return out


# ══════════════════════════════════════════════════════════════
# 页面原型
# ══════════════════════════════════════════════════════════════
def _cover(prs, title: str, meta: list) -> None:
    """原型 1 · 封面：整幅强调色带 + 大标题 + 元信息。"""
    slide = _blank_slide(prs)

    # 顶部色带 + 右下装饰块，形成对角构图
    _rect(slide, 0, 0, SLIDE_W, Inches(0.22), fill=ACCENT)
    _rect(slide, SLIDE_W - Inches(3.1), SLIDE_H - Inches(2.5), Inches(3.1), Inches(2.5),
          fill=SURFACE)
    _rect(slide, SLIDE_W - Inches(3.1), SLIDE_H - Inches(2.5), Pt(4), Inches(2.5), fill=ACCENT)

    # 眉标
    tf = _textbox(slide, MARGIN, Inches(1.72), Inches(9.4), Inches(0.4))
    _para(tf, "三元结构分析报告", 15, ACCENT, True, first=True)

    # 标题（按长度降档，避免溢出）
    size = 38 if len(title) <= 22 else (32 if len(title) <= 34 else 26)
    tf = _textbox(slide, MARGIN, Inches(2.28), Inches(9.4), Inches(2.5))
    _para(tf, title, size, INK, True, first=True, line_spacing=1.22)

    # 分隔短线
    _rect(slide, MARGIN, Inches(4.92), Inches(1.5), Pt(4.5), fill=ACCENT)

    # 元信息（胶囊）
    top = Inches(5.34)
    for line in meta:
        _pill(slide, MARGIN, top, Inches(0.38), line, SURFACE, MUTED, 10.5, False)
        top += Inches(0.48)

    # 右下角装饰：六类利益色点
    dot_x = SLIDE_W - Inches(2.62)
    dot_y = SLIDE_H - Inches(1.92)
    for i, key in enumerate(["gold", "green", "purple", "orange", "blue", "cyan"]):
        _rect(slide, dot_x + i * Inches(0.32), dot_y, Inches(0.2), Inches(0.2),
              fill=C[key], radius=0.5)

    _footer(slide, "基于三元结构理论（三维 × 八阶 × 六类利益）")


def _section_divider(prs, number: str, title: str, accent=ACCENT) -> None:
    """原型 2 · 章节页：大编号 + 标题，留白为主。"""
    slide = _blank_slide(prs)
    _rect(slide, 0, 0, Inches(0.22), SLIDE_H, fill=accent)
    tf = _textbox(slide, MARGIN + Inches(0.3), Inches(2.5), CONTENT_W, Inches(1.5))
    _para(tf, number, 62, accent, True, first=True, line_spacing=1.0)
    tf = _textbox(slide, MARGIN + Inches(0.3), Inches(4.0), CONTENT_W, Inches(1.0))
    _para(tf, title, 30, INK, True, first=True, line_spacing=1.15)
    _rect(slide, MARGIN + Inches(0.3), Inches(5.25), Inches(1.1), Pt(4), fill=accent)


def _overview(prs, report) -> None:
    """原型 3 · 叙述页：大号引导段 + 右侧核心命题卡。"""
    paras = [t for k, t in _items(report, "overview", kinds=("p",))]
    if not paras:
        return
    slide = _blank_slide(prs)
    _title_block(slide, "情况概述", eyebrow="Overview", accent=C["blue"])

    lead = paras[0]
    rest = paras[1:]

    # 左：引导段（大号，左强调边）
    _rect(slide, MARGIN, BODY_TOP, Pt(4), BODY_H, fill=C["blue"])
    tf = _textbox(slide, MARGIN + Inches(0.28), BODY_TOP + Inches(0.04),
                  col_span(7), BODY_H)
    _para(tf, lead, 14.5, INK, first=True, line_spacing=1.55)

    # 右：后续段落做成浅底卡
    if rest:
        # 左 7 列（0-6）为引导段，右 5 列（7-11）为补充卡，合计正好等于版心宽度
        card_l = col_left(7)
        card_w = col_span(5)
        _rect(slide, card_l, BODY_TOP, card_w, BODY_H, fill=SURFACE, radius=0.03)
        tf = _textbox(slide, card_l + Inches(0.3), BODY_TOP + Inches(0.28),
                      card_w - Inches(0.6), BODY_H - Inches(0.56))
        _para(tf, "补充脉络", 11, C["blue"], True, first=True)
        for i, p in enumerate(rest[:2]):
            _para(tf, p, 11.5, MUTED, space_before=8, line_spacing=1.5)

    _footer(slide, "情况概述", 2)


def _conflicts(prs, report) -> None:
    """原型 4 · 冲突卡：主体 A ⟷ 主体 B + 利益徽章 + 说明。"""
    raw = [t for k, t in _items(report, "core_conflicts", li_limit=8, kinds=("li",))]
    if not raw:
        return
    conflicts = [_parse_conflict(r) for r in raw]
    slide = _blank_slide(prs)
    _title_block(slide, "核心冲突点", eyebrow="Conflicts", accent=C["orange"])

    n = len(conflicts)
    cols = 1 if n == 1 else 2
    rows = (n + cols - 1) // cols
    gap = Inches(0.22)
    card_h = Emu(int((BODY_H - gap * (rows - 1)) / rows))
    if card_h > Inches(2.7):
        card_h = Inches(2.7)
    card_w = col_span(GRID // cols) if cols > 1 else CONTENT_W

    for i, cf in enumerate(conflicts[:6]):
        r, c = divmod(i, cols)
        left = col_left(c * (GRID // cols))
        top = Emu(int(BODY_TOP) + r * int(card_h + gap))
        _conflict_card(slide, left, top, card_w, card_h, cf, i + 1)

    _footer(slide, "核心冲突点", 3)


def _conflict_card(slide, left, top, width, height, cf: dict, number: int) -> None:
    """单张冲突卡。"""
    accent = C["orange"]
    _rect(slide, left, top, width, height, fill=CANVAS, line=LINE, radius=0.035)
    _rect(slide, left, top, Pt(4.5), height, fill=accent)
    pad = Inches(0.26)
    inner_l = left + pad
    inner_w = width - pad * 2

    # 编号 + 主体 A ⟷ 主体 B
    _num_circle(slide, inner_l, top + Inches(0.22), Inches(0.34), number, accent)
    y = top + Inches(0.22)
    if cf["a"] or cf["b"]:
        chip_l = inner_l + Inches(0.46)
        chip_h = Inches(0.34)
        chip_y = y + Inches(0.01)
        # 用实际渲染宽度定位，避免胶囊自动缩放后错位或越界
        max_chip = Emu(int((int(inner_w) - int(Inches(0.95))) // 2))
        a_shape = _pill(slide, chip_l, chip_y, chip_h, cf["a"], SURFACE, INK, 10.5, True, max_chip)
        vs_l = chip_l + a_shape.width + Inches(0.1)
        tf = _textbox(slide, vs_l, chip_y + Inches(0.03), Inches(0.34), Inches(0.28))
        _para(tf, "对", 11, accent, True, first=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
        _pill(slide, vs_l + Inches(0.4), chip_y, chip_h, cf["b"], SURFACE, INK, 10.5, True, max_chip)
        y = chip_y + chip_h + Inches(0.16)

    # 利益徽章 + 张力类型
    badge_text = f"{cf['interest']} · {cf['tension']}" if cf["interest"] else cf["tension"]
    if badge_text:
        _pill(slide, inner_l, y, Inches(0.34), badge_text, C["orange"], CANVAS, 9.5, True,
              max_width=inner_w)
        y += Inches(0.5)

    # 说明正文
    detail_h = height - (y - top) - Inches(0.5)
    if detail_h > Inches(0.4) and cf["detail"]:
        tf = _textbox(slide, inner_l, y, inner_w, detail_h)
        size = 10.5 if len(cf["detail"]) > 90 else 11.5
        _para(tf, cf["detail"], size, MUTED, first=True, line_spacing=1.45)

    # 来源
    if cf["source"]:
        tf = _textbox(slide, inner_l, top + height - Inches(0.36), inner_w, Inches(0.26))
        _para(tf, f"来源：{cf['source']}", 8.5, FAINT, first=True, line_spacing=1.0)


def _actors(prs, report) -> None:
    """原型 5 · 主体清单表（v1 因忽略表格而空白）。"""
    tables = _tables(report, "case_portrait")
    if not tables:
        tables = _tables(report, "org_interest_network")
    table = _find_table(tables, ["主体", "角色", "力量"])
    if not table:
        table = tables[0] if tables else None
    if not table:
        # 无表格则退回卡片式列表
        _actors_fallback(prs, report)
        return

    slide = _blank_slide(prs)
    _title_block(slide, "利益主体识别", eyebrow="Actors", accent=C["blue"])
    rows = table["rows"]
    _render_table(slide, rows, MARGIN, BODY_TOP, CONTENT_W, BODY_H,
                  header_fill=C["blue"], caption=table["caption"])
    _footer(slide, "利益主体识别", 4)


def _actors_fallback(prs, report) -> None:
    """无表格时：主体卡片网格。"""
    items = _items(report, "case_portrait", li_limit=8)
    if not items:
        items = _items(report, "org_interest_network", li_limit=8)
    if not items:
        return
    slide = _blank_slide(prs)
    _title_block(slide, "利益主体识别", eyebrow="Actors", accent=C["blue"])
    tf = _textbox(slide, MARGIN, BODY_TOP, CONTENT_W, BODY_H)
    for i, (kind, text) in enumerate(items[:10]):
        _para(tf, text, 13, INK if kind != "sub" else C["blue"], kind == "sub",
              first=(i == 0), space_after=9, line_spacing=1.4)
    _footer(slide, "利益主体识别", 4)


def _interest_matrix(prs, report) -> None:
    """原型 6 · 六类利益分布矩阵（理论专属视觉）。"""
    tables = _tables(report, "case_portrait")
    if not tables:
        tables = _tables(report, "org_interest_network")
    table = _find_table(tables, ["物质", "安全", "政治", "身份", "制度", "公共"])
    if not table:
        return
    rows = table["rows"]
    slide = _blank_slide(prs)
    _title_block(slide, "六类利益分布", eyebrow="Interests", accent=C["gold"])

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    avail_h = min(BODY_H, Inches(0.62) * n_rows)
    top = BODY_TOP + Inches(0.1)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, MARGIN, top, CONTENT_W, avail_h
    )
    tbl = table_shape.table
    tbl.first_row = True
    tbl.horz_banding = False

    try:
        col_widths = [Inches(1.95)] + [int((CONTENT_W - Inches(1.95)) / (n_cols - 1))] * (n_cols - 1)
        for i, w in enumerate(col_widths[:n_cols]):
            tbl.columns[i].width = Emu(int(w))
    except Exception:  # noqa: BLE001
        pass

    for r_i, row in enumerate(rows):
        for c_i in range(n_cols):
            cell = tbl.cell(r_i, c_i)
            text = row[c_i] if c_i < len(row) else ""
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _interest_color(text) if c_i else C["gold"]
                _cell_text(cell, text, 10.5, CANVAS, True, PP_ALIGN.CENTER)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE if r_i % 2 else CANVAS
                color = INK if c_i == 0 else _level_color(text)
                bold = c_i == 0
                _cell_text(cell, text, 10, color, bold,
                           PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER)

    _footer(slide, "六类利益分布", 5)


def _relations(prs, report) -> None:
    """原型 7 · 主体间关系网络（双列键值卡）。"""
    tables = _tables(report, "case_portrait")
    if not tables:
        tables = _tables(report, "org_interest_network")
    table = _find_table(tables, ["关系维度", "维度"])
    if not table:
        return
    rows = table["rows"]
    body = [r for r in rows[1:] if r and r[0]]
    if not body:
        return

    slide = _blank_slide(prs)
    _title_block(slide, "主体间关系网络", eyebrow="Relations", accent=C["purple"])

    n = len(body)
    cols = 2 if n > 2 else 1
    rows_n = (n + cols - 1) // cols
    gap = Inches(0.18)
    card_h = Emu(int((BODY_H - gap * (rows_n - 1)) / rows_n))
    if card_h > Inches(1.4):
        card_h = Inches(1.4)

    for i, row in enumerate(body[:6]):
        r, c = divmod(i, cols)
        left = col_left(c * (GRID // cols))
        top = Emu(int(BODY_TOP) + r * int(card_h + gap))
        w = col_span(GRID // cols) if cols > 1 else CONTENT_W
        key = row[0]
        val = row[1] if len(row) > 1 else ""
        note = row[2] if len(row) > 2 else ""
        _rect(slide, left, top, w, card_h, fill=SURFACE, radius=0.04)
        _pill(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.34), key,
              C["purple"], CANVAS, 10, True, max_width=w - Inches(0.4))
        tf = _textbox(slide, left + Inches(0.2), top + Inches(0.66),
                      w - Inches(0.4), card_h - Inches(0.8))
        cur = val
        if note:
            cur = f"{val} · {note}" if val else note
        _para(tf, cur, 10.5, MUTED, first=True, line_spacing=1.4)

    _footer(slide, "主体间关系网络", 6)


def _render_table(slide, rows, left, top, width, height, header_fill=ACCENT,
                  caption=None) -> None:
    """通用表格渲染：强调色表头 + 斑马纹。"""
    if not rows:
        return
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    if caption:
        tf = _textbox(slide, left, top - Inches(0.34), width, Inches(0.28))
        _para(tf, caption, 11, header_fill, True, first=True)
    avail_h = min(height, Inches(0.52) * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, avail_h)
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False

    # 列宽：首列略宽
    if n_cols > 1:
        first = Inches(2.55)
        rest = Emu(int((width - first) / (n_cols - 1)))
        try:
            tbl.columns[0].width = Emu(int(first))
            for i in range(1, n_cols):
                tbl.columns[i].width = rest
        except Exception:  # noqa: BLE001
            pass

    for r_i, row in enumerate(rows):
        for c_i in range(n_cols):
            cell = tbl.cell(r_i, c_i)
            text = row[c_i] if c_i < len(row) else ""
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r_i == 0:
                cell.fill.fore_color.rgb = header_fill
                _cell_text(cell, text, 11, CANVAS, True, PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER)
            else:
                cell.fill.fore_color.rgb = SURFACE if r_i % 2 else CANVAS
                _cell_text(cell, text, 10,
                           INK if c_i == 0 else MUTED, c_i == 0,
                           PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.CENTER)


def _cell_text(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.25
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    if text:
        run = p.add_run()
        run.text = text
        _apply_font(run, size, color, bold)


def _diagram_slide(prs, png_path: str, title: str) -> None:
    """原型 8 · 全幅关系图：关系图独占几乎整页，标题作为小角标。"""
    slide = _blank_slide(prs)
    # 小标题叠在左上角，不占用大量纵向空间
    tf = _textbox(slide, MARGIN, Inches(0.34), CONTENT_W, Inches(0.46))
    _para(tf, title or "利益关系网络", 11, MUTED, True, first=True)
    _rect(slide, MARGIN, Inches(0.82), Inches(0.8), Pt(3), fill=C["purple"])

    render_path = png_path
    try:
        from PIL import Image

        img = Image.open(png_path)
        if img.mode == "RGBA":
            bbox = img.getbbox()
            if bbox and bbox != (0, 0, img.width, img.height):
                img = img.crop(bbox)
        tmp = os.path.join(os.path.dirname(png_path), "_cropped_diagram.png")
        img.save(tmp)
        render_path = tmp
        # 关键：PIL 返回像素，必须按 96 DPI 转成英寸/EMU 后再算比例
        img_w = Inches(img.width / 96.0)
        img_h = Inches(img.height / 96.0)
    except Exception:  # noqa: BLE001
        render_path = png_path
        img_w, img_h = Inches(12.5), Inches(8.33)

    # 上下留少量呼吸边距，其余全部给图
    margin_v = Inches(0.42)
    avail_w = SLIDE_W - 2 * MARGIN
    avail_h = SLIDE_H - Inches(1.05) - 2 * margin_v
    scale = min(float(avail_w) / float(img_w), float(avail_h) / float(img_h), 1.0)
    disp_w = Emu(int(img_w * scale))
    disp_h = Emu(int(img_h * scale))
    left = int((SLIDE_W - disp_w) / 2)
    top = int(Inches(1.05) + margin_v + (avail_h - disp_h) / 2)
    _rect(slide, left - Inches(0.1), top - Inches(0.1),
          disp_w + Inches(0.2), disp_h + Inches(0.2), fill=SURFACE, radius=0.015)
    slide.shapes.add_picture(render_path, left, top, disp_w, disp_h)
    _footer(slide, "利益关系网络图")


def _analysis(prs, report) -> None:
    """原型 9 · 三元结构分析：三栏（主体 / 利益 / 关系）。"""
    items = _items(report, "analysis_body", li_limit=12)
    if not items:
        return
    paras = [t for k, t in items if k == "p"]
    subs = [t for k, t in items if k == "sub"]
    lis = [t for k, t in items if k == "li"]

    slide = _blank_slide(prs)
    _title_block(slide, "三元结构分析", eyebrow="Analysis", accent=C["green"])

    # 三栏：优先子标题分栏，否则按 主体/利益/关系 切分内容
    col_defs = [
        ("主体", C["gold"], subs[:3] or paras[:1] or lis[:3]),
        ("利益", C["green"], subs[3:6] or paras[1:2] or lis[3:6]),
        ("关系", C["blue"], subs[6:9] or paras[2:3] or lis[6:9]),
    ]
    span = 4
    for i, (name, color, chunk) in enumerate(col_defs):
        left = col_left(i * span)
        w = col_span(span)
        _rect(slide, left, BODY_TOP, w, BODY_H, fill=SURFACE, radius=0.03)
        _rect(slide, left, BODY_TOP, w, Inches(0.52), fill=color, radius=None)
        tf = _textbox(slide, left, BODY_TOP, w, Inches(0.52))
        _para(tf, name, 13, CANVAS, True, first=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
        tf = _textbox(slide, left + Inches(0.24), BODY_TOP + Inches(0.78),
                      w - Inches(0.48), BODY_H - Inches(1.0))
        if chunk:
            for j, txt in enumerate(chunk[:4]):
                _para(tf, txt, 11, INK, first=(j == 0), space_after=10, line_spacing=1.45)
        else:
            _para(tf, "（本节未提供）", 10.5, FAINT, first=True)

    _footer(slide, "三元结构分析", 7)


def _conclusion(prs, report) -> None:
    """原型 10 · 结论：大号引言式陈述。"""
    items = [t for k, t in _items(report, "conclusion", li_limit=8, kinds=("p", "li"))]
    if not items:
        return
    slide = _blank_slide(prs)
    _title_block(slide, "核心结论", eyebrow="Conclusion", accent=C["gold"])

    # 主结论：大号引言
    lead = items[0]
    size = 22 if len(lead) <= 60 else (18 if len(lead) <= 110 else 15)
    _rect(slide, MARGIN, BODY_TOP, CONTENT_W, Inches(2.5), fill=SURFACE, radius=0.03)
    _rect(slide, MARGIN, BODY_TOP, Pt(5), Inches(2.5), fill=C["gold"])
    tf = _textbox(slide, MARGIN + Inches(0.4), BODY_TOP + Inches(0.42),
                  CONTENT_W - Inches(0.8), Inches(1.7), anchor=MSO_ANCHOR.MIDDLE)
    _para(tf, lead, size, INK, True, first=True, align=PP_ALIGN.LEFT, line_spacing=1.42)

    # 其余结论：编号列表
    rest = items[1:5]
    if rest:
        y = BODY_TOP + Inches(2.82)
        for i, txt in enumerate(rest):
            _num_circle(slide, MARGIN, y, Inches(0.3), i + 1, C["gold"])
            tf = _textbox(slide, MARGIN + Inches(0.46), y + Inches(0.03),
                          CONTENT_W - Inches(0.46), Inches(0.72))
            _para(tf, txt, 12, INK, first=True, line_spacing=1.42)
            y += Inches(0.56)

    _footer(slide, "核心结论", 8)


def _recommendations(prs, report) -> None:
    """原型 11 · 行动建议：编号圆 + 双列。"""
    items = [t for k, t in _items(report, "recommendations", li_limit=12, kinds=("li",))]
    if not items:
        return
    slide = _blank_slide(prs)
    _title_block(slide, "行动建议", eyebrow="Actions", accent=C["cyan"])

    n = len(items)
    cols = 2 if n > 3 else 1
    rows_n = (n + cols - 1) // cols
    gap = Inches(0.2)
    row_h = Emu(int((BODY_H - gap * (rows_n - 1)) / rows_n))
    if row_h > Inches(1.25):
        row_h = Inches(1.25)

    for i, txt in enumerate(items[:8]):
        r, c = divmod(i, cols)
        left = col_left(c * (GRID // cols))
        top = Emu(int(BODY_TOP) + r * int(row_h + gap))
        w = col_span(GRID // cols) if cols > 1 else CONTENT_W
        _rect(slide, left, top, w, row_h, fill=SURFACE, radius=0.06)
        _num_circle(slide, left + Inches(0.22), top + Inches(0.24), Inches(0.34), i + 1, C["cyan"])
        tf = _textbox(slide, left + Inches(0.72), top + Inches(0.2),
                      w - Inches(0.98), row_h - Inches(0.36), anchor=MSO_ANCHOR.MIDDLE)
        size = 11.5 if len(txt) > 70 else 12.5
        _para(tf, txt, size, INK, first=True, line_spacing=1.4)

    _footer(slide, "行动建议", 9)


def _appendix(prs, report) -> None:
    """原型 12 · 附录来源：双栏小字列表。"""
    srcs = _sources(report)
    slide = _blank_slide(prs)
    _title_block(slide, "附录 · 数据来源", eyebrow="Sources", accent=C["blue"])
    tf_left = _textbox(slide, MARGIN, BODY_TOP, col_span(6), BODY_H)
    tf_right = _textbox(slide, col_left(6), BODY_TOP, col_span(6), BODY_H)
    if not srcs:
        _para(tf_left, "（本报告未引用外部来源）", 11.5, FAINT, first=True)
        _footer(slide, "附录", 10)
        return
    half = (len(srcs) + 1) // 2
    for i, (name, url) in enumerate(srcs):
        tf = tf_left if i < half else tf_right
        idx = i if i < half else i - half
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.22
        run = p.add_run()
        run.text = f"{i + 1:02d}  "
        _apply_font(run, 9.5, C["blue"], True)
        run2 = p.add_run()
        run2.text = name[:30] + ("…" if len(name) > 30 else "")
        _apply_font(run2, 10.5, INK)
        if url:
            run3 = p.add_run()
            run3.text = f"\n     {url[:48]}"
            _apply_font(run3, 8, FAINT)
    _footer(slide, f"共 {len(srcs)} 项来源", 10)


# ══════════════════════════════════════════════════════════════
# 编排
# ══════════════════════════════════════════════════════════════
def render_report_pptx(
    title: str,
    markdown: str,
    output_path: str,
    diagram_pngs: list | None = None,
) -> str:
    """把报告 Markdown 渲染为 .pptx；diagram_pngs 可选 [{seq,png,title}]。"""
    _ensure_engine_path()
    from parser import parse_report

    report = parse_report(markdown)
    if not report.title and title:
        report.title = title

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    meta = [
        f"分析类型：{_analysis_type_label(report)}",
        f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}",
    ]
    _cover(prs, report.title or title, meta)
    _overview(prs, report)
    _conflicts(prs, report)
    _actors(prs, report)
    _interest_matrix(prs, report)
    _relations(prs, report)
    for diag in diagram_pngs or []:
        _diagram_slide(prs, diag["png"], diag.get("title") or "利益关系网络")
    _analysis(prs, report)
    _conclusion(prs, report)
    _recommendations(prs, report)
    _appendix(prs, report)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _analysis_type_label(report) -> str:
    """按章节哨兵推断报告类型（policy/event/org/opinion/combo）。"""
    ids = set(report.sections.keys())
    if {"policy_portrait", "policy_weight"} & ids:
        return "政策分析"
    if {"org_structure", "org_survival", "org_interest_network"} & ids:
        return "组织诊断"
    if {"opinion_event", "opinion_actors", "opinion_narrative"} & ids:
        return "舆情分析"
    if {"case_portrait", "case_flows"} & ids:
        return "事件分析"
    return "结构分析"


def export_report_pptx(
    title: str,
    markdown: str,
    output_dir: str | None = None,
    slug: str | None = None,
) -> dict:
    """导出 .pptx：在 output_dir/slug 下生成 diagram PNG + pptx。"""
    base = output_dir or settings.generated_dir
    safe = re.sub(r"[\\/:*?\"<>|]", "-", (title or "未命名报告").strip()) or "未命名报告"
    folder = os.path.join(base, f"{safe}_{datetime.now().strftime('%Y%m%d_%H%M%S')}")
    if slug:
        folder = os.path.join(base, slug)
    os.makedirs(folder, exist_ok=True)

    _ensure_engine_path()
    from parser import parse_report

    report = parse_report(markdown)
    diagrams = _diagram_pngs(report, folder)
    pptx_path = os.path.join(folder, f"{safe}.pptx")
    render_report_pptx(report.title or title, markdown, pptx_path, diagram_pngs=diagrams)
    return {
        "pptx": pptx_path,
        "folder": folder,
        "diagram_count": len(diagrams),
        "diagrams": [
            {"png": d["png"], "title": d.get("title") or f"图 {d['seq']}"} for d in diagrams
        ],
    }
